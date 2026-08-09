from __future__ import annotations

import re
from dataclasses import dataclass

from lxml import etree
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from .colors import parse_color
from .coordinates import CoordinateMapper
from .fonts import (
    baseline_to_top,
    estimate_text_width,
    is_bold,
    line_box_height,
    primary_font,
)
from .pptx_shapes import NativeShapeUnsupported
from .pptx_utils import set_font_color
from .svg_models import RenderElement
from .svg_parser import element_styles, local_name, parse_length, parse_number

XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


@dataclass(frozen=True, slots=True)
class TextRunSpec:
    text: str
    styles: dict[str, str]


@dataclass(frozen=True, slots=True)
class TextLineSpec:
    x: float
    y: float
    anchor: str
    runs: tuple[TextRunSpec, ...]


def extract_text_lines(item: RenderElement) -> list[TextLineSpec]:
    element = item.element
    base_x = parse_number(element.get("x"), 0.0)
    base_y = parse_number(element.get("y"), 0.0)
    current_x, current_y = base_x, base_y
    anchor = item.styles.get("text-anchor", element.get("text-anchor", "start"))
    lines: list[TextLineSpec] = []
    runs: list[TextRunSpec] = []

    leading_text = _normalized_text_node(element.text, element)
    if leading_text:
        runs.append(TextRunSpec(leading_text, dict(item.styles)))

    for child in element:
        if not isinstance(child.tag, str) or local_name(child.tag) != "tspan":
            continue
        child_styles = element_styles(child, item.styles)
        next_x = parse_number(child.get("x"), current_x) + parse_number(child.get("dx"), 0.0)
        next_y = parse_number(child.get("y"), current_y) + parse_number(child.get("dy"), 0.0)
        begins_line = any(child.get(name) is not None for name in ("x", "y", "dy"))
        if begins_line and runs:
            lines.append(TextLineSpec(current_x, current_y, anchor, tuple(runs)))
            runs = []
        current_x, current_y = next_x, next_y
        text = _normalized_text_node("".join(child.itertext()), child)
        if text:
            runs.append(TextRunSpec(text, child_styles))
        tail = _normalized_text_node(child.tail, element)
        if tail:
            runs.append(TextRunSpec(tail, dict(item.styles)))

    if runs:
        lines.append(TextLineSpec(current_x, current_y, anchor, tuple(runs)))
    return lines


def add_native_text(
    slide: object,
    item: RenderElement,
    mapper: CoordinateMapper,
) -> list[object]:
    if not item.transform.is_axis_aligned:
        raise NativeShapeUnsupported("rotated or skewed text")
    dominant_baseline = item.styles.get("dominant-baseline", "").strip().lower()
    if dominant_baseline not in {"", "auto", "alphabetic"}:
        raise NativeShapeUnsupported(
            f"dominant-baseline={dominant_baseline!r} is not supported by native text"
        )
    if any("url(" in item.styles.get(key, "").lower() for key in ("fill", "stroke")):
        raise NativeShapeUnsupported("gradient text")

    shapes: list[object] = []
    for line_number, line in enumerate(extract_text_lines(item), start=1):
        if not any(run.text.strip() for run in line.runs):
            continue
        max_font_size = max(_font_size(run.styles) for run in line.runs)
        scaled_font_size = max_font_size * abs(item.transform.d)
        raw_width = sum(
            estimate_text_width(
                run.text,
                _font_size(run.styles),
                bold=is_bold(run.styles.get("font-weight")),
            )
            for run in line.runs
        )
        width = max(1.0, raw_width * abs(item.transform.a))
        baseline_x, baseline_y = item.transform.apply(line.x, line.y)
        left = baseline_x
        if line.anchor == "middle":
            left -= width / 2
        elif line.anchor == "end":
            left -= width
        top = baseline_to_top(baseline_y, scaled_font_size)
        height = line_box_height(scaled_font_size)

        shape = slide.shapes.add_textbox(
            mapper.x(left),
            mapper.y(top),
            mapper.width(width + scaled_font_size * 0.08),
            mapper.height(height),
        )
        shape.name = item.element.get("id") or f"SVG text {line_number}"
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.word_wrap = False
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = _paragraph_alignment(line.anchor)
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        for run_spec in line.runs:
            run = paragraph.add_run()
            run.text = run_spec.text
            font_size = _font_size(run_spec.styles) * abs(item.transform.d)
            run.font.name = primary_font(run_spec.styles.get("font-family"))
            run.font.size = Pt(mapper.font_size_points(font_size))
            run.font.bold = is_bold(run_spec.styles.get("font-weight"))
            run.font.italic = run_spec.styles.get("font-style", "").lower() in {
                "italic",
                "oblique",
            }
            color = parse_color(
                run_spec.styles.get("fill", "#000000"),
                opacity=item.opacity * _opacity(run_spec.styles.get("fill-opacity")),
            )
            if color is not None:
                set_font_color(run, color)
        shapes.append(shape)
    return shapes


def _font_size(styles: dict[str, str]) -> float:
    value = parse_length(styles.get("font-size"))
    return value if value is not None else 16.0


def _normalized_text_node(value: str | None, scope: etree._Element) -> str:
    if value is None:
        return ""
    if _preserves_whitespace(scope):
        return value
    if not value.strip():
        return ""
    return re.sub(r"\s+", " ", value)


def _preserves_whitespace(element: etree._Element) -> bool:
    current: etree._Element | None = element
    while current is not None:
        value = current.get(XML_SPACE)
        if value is not None:
            return value.strip().lower() == "preserve"
        current = current.getparent()
    return False


def _opacity(value: str | None) -> float:
    try:
        return max(0.0, min(1.0, float(value))) if value is not None else 1.0
    except ValueError:
        return 1.0


def _paragraph_alignment(anchor: str) -> PP_ALIGN:
    if anchor == "middle":
        return PP_ALIGN.CENTER
    if anchor == "end":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT
