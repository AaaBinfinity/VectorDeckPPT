from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from .colors import SvgColor
from .path_safety import require_suffix

SVG_BLIP_NAMESPACE = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_BLIP_EXTENSION_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
RELATIONSHIP_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


@dataclass(slots=True)
class SlideCompilationReport:
    source: str
    native: int = 0
    freeform: int = 0
    embedded_svg: int = 0
    failed: int = 0
    warnings: list[str] = field(default_factory=list)


@dataclass(slots=True)
class CompilationReport:
    source: str
    output: str
    slides: list[SlideCompilationReport] = field(default_factory=list)

    @property
    def native(self) -> int:
        return sum(slide.native for slide in self.slides)

    @property
    def freeform(self) -> int:
        return sum(slide.freeform for slide in self.slides)

    @property
    def embedded_svg(self) -> int:
        return sum(slide.embedded_svg for slide in self.slides)

    @property
    def failed(self) -> int:
        return sum(slide.failed for slide in self.slides)

    @property
    def valid(self) -> bool:
        return bool(self.slides) and self.failed == 0

    def to_dict(self) -> dict[str, object]:
        return {
            "valid": self.valid,
            "source": self.source,
            "output": self.output,
            "slide_count": len(self.slides),
            "native": self.native,
            "freeform": self.freeform,
            "embedded_svg": self.embedded_svg,
            "failed": self.failed,
            "slides": [asdict(slide) for slide in self.slides],
        }

    def write_json(self, path: str | Path) -> Path:
        output = require_suffix(path, ".json", label="Compilation report")
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(
            json.dumps(self.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8"
        )
        return output


def natural_sort_key(path: Path) -> list[int | str]:
    return [
        int(chunk) if chunk.isdigit() else chunk.casefold()
        for chunk in re.split(r"(\d+)", path.name)
    ]


def set_fill(fill_format: object, color: SvgColor | None) -> None:
    if color is None:
        fill_format.background()
        return
    fill_format.solid()
    fill_format.fore_color.rgb = RGBColor(*color.rgb)
    _set_alpha(fill_format._xPr, color.alpha)


def set_line(
    line_format: object,
    color: SvgColor | None,
    width_emu: int,
    *,
    cap: str | None = None,
    join: str | None = None,
) -> None:
    if color is None:
        line_format.fill.background()
        return
    line_format.fill.solid()
    line_format.fill.fore_color.rgb = RGBColor(*color.rgb)
    line_format.width = width_emu
    _set_alpha(line_format.fill._xPr, color.alpha)
    line_properties = line_format._get_or_add_ln()
    _set_line_cap(line_properties, cap)
    _set_line_join(line_properties, join)


def _set_line_cap(line_properties: object, value: str | None) -> None:
    if value is None or not value.strip():
        return
    caps = {"butt": "flat", "round": "rnd", "square": "sq"}
    normalized = value.strip().lower()
    if normalized not in caps:
        raise ValueError(f"Unsupported SVG stroke-linecap: {value!r}")
    line_properties.set("cap", caps[normalized])


def _set_line_join(line_properties: object, value: str | None) -> None:
    if value is None or not value.strip():
        return
    joins = {"miter": "a:miter", "round": "a:round", "bevel": "a:bevel"}
    normalized = value.strip().lower()
    if normalized not in joins:
        raise ValueError(f"Unsupported SVG stroke-linejoin: {value!r}")
    for tag in joins.values():
        for existing in line_properties.xpath(f"./{tag}"):
            line_properties.remove(existing)
    line_properties.append(OxmlElement(joins[normalized]))


def set_font_color(run: object, color: SvgColor) -> None:
    run.font.color.rgb = RGBColor(*color.rgb)
    run_properties = run._r.get_or_add_rPr()
    _set_alpha(run_properties, color.alpha)


def clear_shape_theme_style(shape: object) -> None:
    """Remove PowerPoint theme effects that are not present in the source SVG."""

    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


def _set_alpha(parent: object, alpha: float) -> None:
    color_nodes = parent.xpath(".//a:solidFill/*")
    if not color_nodes:
        return
    color_node = color_nodes[-1]
    for existing in color_node.xpath("./a:alpha"):
        color_node.remove(existing)
    if alpha >= 0.9999:
        return
    alpha_node = OxmlElement("a:alpha")
    alpha_node.set("val", str(round(max(0.0, min(1.0, alpha)) * 100_000)))
    color_node.append(alpha_node)


def add_svg_picture(
    slide: object,
    svg_data: bytes,
    preview_png: bytes,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    name: str,
) -> object:
    """Add a PNG-compatible picture with an Office SVG extension relationship."""

    picture = slide.shapes.add_picture(BytesIO(preview_png), left, top, width, height)
    picture.name = name

    package = slide.part.package
    svg_partname = package.next_image_partname("svg")
    svg_part = Part(svg_partname, "image/svg+xml", package, svg_data)
    svg_relationship_id = slide.part.relate_to(svg_part, RT.IMAGE)

    blip = picture._element.blipFill.blip
    extension_list = blip.find("{http://schemas.openxmlformats.org/drawingml/2006/main}extLst")
    if extension_list is None:
        extension_list = OxmlElement("a:extLst")
        blip.append(extension_list)
    extension = OxmlElement("a:ext")
    extension.set("uri", SVG_BLIP_EXTENSION_URI)
    svg_blip = etree.SubElement(
        extension,
        f"{{{SVG_BLIP_NAMESPACE}}}svgBlip",
        nsmap={"asvg": SVG_BLIP_NAMESPACE},
    )
    svg_blip.set(f"{{{RELATIONSHIP_NAMESPACE}}}embed", svg_relationship_id)
    extension_list.append(extension)
    return picture
