from __future__ import annotations

import posixpath
import re
from dataclasses import dataclass, field
from pathlib import Path, PurePosixPath
from typing import Any
from zipfile import BadZipFile, ZipFile

from lxml import etree
from pptx import Presentation

RELATIONSHIP_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"
SLIDE_RE = re.compile(r"^ppt/slides/slide\d+\.xml$")
REQUIRED_PARTS = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}


@dataclass(slots=True)
class PptxValidationResult:
    source: Path
    slides: int = 0
    media: int = 0
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def valid(self) -> bool:
        return not self.errors

    def to_dict(self) -> dict[str, Any]:
        return {
            "valid": self.valid,
            "source": str(self.source),
            "slides": self.slides,
            "media": self.media,
            "errors": self.errors,
            "warnings": self.warnings,
        }


def validate_pptx(path: str | Path) -> PptxValidationResult:
    source = Path(path).expanduser().resolve()
    result = PptxValidationResult(source=source)
    if not source.is_file():
        result.errors.append(f"PPTX file does not exist: {source}")
        return result
    if source.suffix.lower() != ".pptx":
        result.warnings.append("File extension is not .pptx")

    try:
        with ZipFile(source) as archive:
            names = set(archive.namelist())
            corrupt_member = archive.testzip()
            if corrupt_member:
                result.errors.append(f"Corrupt ZIP member: {corrupt_member}")
            missing = sorted(REQUIRED_PARTS - names)
            if missing:
                result.errors.append(f"Missing required PPTX parts: {', '.join(missing)}")
            result.slides = sum(bool(SLIDE_RE.match(name)) for name in names)
            result.media = sum(
                name.startswith("ppt/media/") and not name.endswith("/") for name in names
            )
            if result.slides == 0:
                result.errors.append("PPTX contains no slide parts")
            _validate_relationships(archive, names, result)
            _validate_xml_parts(archive, names, result)
    except BadZipFile as exc:
        result.errors.append(f"Invalid ZIP/PPTX container: {exc}")
        return result
    except OSError as exc:
        result.errors.append(f"Unable to read PPTX: {exc}")
        return result

    if not result.errors:
        try:
            presentation = Presentation(source)
            if len(presentation.slides) != result.slides:
                result.errors.append(
                    "Presentation slide count does not match slide parts: "
                    f"{len(presentation.slides)} != {result.slides}"
                )
        except Exception as exc:
            result.errors.append(f"python-pptx could not open the deck: {exc}")
    return result


def _validate_relationships(
    archive: ZipFile,
    names: set[str],
    result: PptxValidationResult,
) -> None:
    for relationship_name in sorted(name for name in names if name.endswith(".rels")):
        try:
            root = etree.fromstring(archive.read(relationship_name))
        except etree.XMLSyntaxError as exc:
            result.errors.append(f"Malformed relationship XML {relationship_name}: {exc}")
            continue
        source_part = _source_part_for_relationship(relationship_name)
        base = posixpath.dirname(source_part)
        for relationship in root.findall(f"{{{RELATIONSHIP_NAMESPACE}}}Relationship"):
            if relationship.get("TargetMode") == "External":
                result.warnings.append(
                    f"External relationship in {relationship_name}: {relationship.get('Target')}"
                )
                continue
            target = relationship.get("Target")
            if not target:
                result.errors.append(f"Relationship without Target in {relationship_name}")
                continue
            resolved = (
                target.lstrip("/")
                if target.startswith("/")
                else posixpath.normpath(posixpath.join(base, target))
            )
            if resolved not in names:
                result.errors.append(
                    f"Broken relationship in {relationship_name}: {target} -> {resolved}"
                )


def _validate_xml_parts(
    archive: ZipFile,
    names: set[str],
    result: PptxValidationResult,
) -> None:
    important_parts = [
        "ppt/presentation.xml",
        *sorted(name for name in names if SLIDE_RE.match(name)),
    ]
    for name in important_parts:
        if name not in names:
            continue
        try:
            etree.fromstring(archive.read(name))
        except etree.XMLSyntaxError as exc:
            result.errors.append(f"Malformed XML part {name}: {exc}")


def _source_part_for_relationship(relationship_name: str) -> str:
    path = PurePosixPath(relationship_name)
    if relationship_name == "_rels/.rels":
        return ""
    parent = path.parent
    if parent.name != "_rels":
        return str(path)
    source_name = path.name.removesuffix(".rels")
    return str(parent.parent / source_name)
