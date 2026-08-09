from __future__ import annotations

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE

from .colors import parse_color
from .coordinates import CoordinateMapper
from .pptx_utils import clear_shape_theme_style, set_fill, set_line
from .svg_models import Matrix, RenderElement, SvgDocument
from .svg_parser import parse_length, parse_points


class NativeShapeUnsupported(ValueError):
    """Raised when a visible SVG element needs embedded-SVG fallback."""


def add_native_shape(
    slide: object,
    item: RenderElement,
    document: SvgDocument,
    mapper: CoordinateMapper,
) -> object:
    if not item.transform.is_axis_aligned:
        raise NativeShapeUnsupported("rotated or skewed shape")
    if _requires_svg_fallback(item):
        raise NativeShapeUnsupported(
            "paint server, clipping, marker, or dashed stroke is not supported natively"
        )

    if item.tag == "rect":
        shape = _add_rect(slide, item, document, mapper)
    elif item.tag in {"circle", "ellipse"}:
        shape = _add_ellipse(slide, item, document, mapper)
    elif item.tag == "line":
        shape = _add_line(slide, item, document, mapper)
    elif item.tag in {"polyline", "polygon"}:
        shape = _add_freeform(slide, item, mapper)
    else:
        raise NativeShapeUnsupported(f"unsupported native shape <{item.tag}>")
    clear_shape_theme_style(shape)
    shape.name = item.element.get("id") or f"SVG {item.tag}"
    return shape


def _add_rect(
    slide: object,
    item: RenderElement,
    document: SvgDocument,
    mapper: CoordinateMapper,
) -> object:
    element = item.element
    x = _length(element, "x", document, default=0)
    y = _length(element, "y", document, default=0, horizontal=False)
    width = _length(element, "width", document)
    height = _length(element, "height", document, horizontal=False)
    left, top, right, bottom = _transformed_box(item.transform, x, y, width, height)
    if right <= left or bottom <= top:
        raise NativeShapeUnsupported("zero-size rectangle")

    rx = _length(element, "rx", document, default=0)
    ry = _length(element, "ry", document, default=rx, horizontal=False)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if max(rx, ry) > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        mapper.x(left),
        mapper.y(top),
        mapper.width(right - left),
        mapper.height(bottom - top),
    )
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE and shape.adjustments:
        radius = max(rx * abs(item.transform.a), ry * abs(item.transform.d))
        shape.adjustments[0] = min(0.5, radius / min(right - left, bottom - top))
    _apply_shape_style(shape, item, mapper)
    return shape


def _add_ellipse(
    slide: object,
    item: RenderElement,
    document: SvgDocument,
    mapper: CoordinateMapper,
) -> object:
    element = item.element
    cx = _length(element, "cx", document, default=0)
    cy = _length(element, "cy", document, default=0, horizontal=False)
    if item.tag == "circle":
        rx = ry = _length(element, "r", document)
    else:
        rx = _length(element, "rx", document)
        ry = _length(element, "ry", document, horizontal=False)
    left, top, right, bottom = _transformed_box(
        item.transform,
        cx - rx,
        cy - ry,
        rx * 2,
        ry * 2,
    )
    if right <= left or bottom <= top:
        raise NativeShapeUnsupported("zero-size ellipse")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        mapper.x(left),
        mapper.y(top),
        mapper.width(right - left),
        mapper.height(bottom - top),
    )
    _apply_shape_style(shape, item, mapper)
    return shape


def _add_line(
    slide: object,
    item: RenderElement,
    document: SvgDocument,
    mapper: CoordinateMapper,
) -> object:
    element = item.element
    x1 = _length(element, "x1", document, default=0)
    y1 = _length(element, "y1", document, default=0, horizontal=False)
    x2 = _length(element, "x2", document, default=0)
    y2 = _length(element, "y2", document, default=0, horizontal=False)
    start_x, start_y = item.transform.apply(x1, y1)
    end_x, end_y = item.transform.apply(x2, y2)
    shape = slide.shapes.add_connector(
        1,
        mapper.x(start_x),
        mapper.y(start_y),
        mapper.x(end_x),
        mapper.y(end_y),
    )
    stroke = parse_color(
        item.styles.get("stroke"),
        opacity=item.opacity * _opacity(item.styles.get("stroke-opacity")),
    )
    stroke_width = parse_length(item.styles.get("stroke-width")) or 1.0
    set_line(
        shape.line,
        stroke,
        mapper.width(stroke_width * abs(item.transform.a)),
        cap=item.styles.get("stroke-linecap"),
        join=item.styles.get("stroke-linejoin"),
    )
    return shape


def _add_freeform(
    slide: object,
    item: RenderElement,
    mapper: CoordinateMapper,
) -> object:
    points = parse_points(item.element.get("points"))
    minimum_points = 3 if item.tag == "polygon" else 2
    if len(points) < minimum_points:
        raise NativeShapeUnsupported(
            f"<{item.tag}> requires at least {minimum_points} coordinate pairs"
        )
    transformed = [item.transform.apply(x, y) for x, y in points]
    local_points = [
        (x - mapper.view_x, y - mapper.view_y)
        for x, y in transformed
    ]
    builder = slide.shapes.build_freeform(
        start_x=local_points[0][0],
        start_y=local_points[0][1],
        scale=(mapper.width(1.0), mapper.height(1.0)),
    )
    builder.add_line_segments(local_points[1:], close=item.tag == "polygon")
    shape = builder.convert_to_shape()
    _apply_shape_style(shape, item, mapper)
    return shape


def _apply_shape_style(shape: object, item: RenderElement, mapper: CoordinateMapper) -> None:
    fill = parse_color(
        item.styles.get("fill", "#000000"),
        opacity=item.opacity * _opacity(item.styles.get("fill-opacity")),
    )
    stroke = parse_color(
        item.styles.get("stroke"),
        opacity=item.opacity * _opacity(item.styles.get("stroke-opacity")),
    )
    stroke_width = parse_length(item.styles.get("stroke-width")) or 1.0
    set_fill(shape.fill, fill)
    set_line(
        shape.line,
        stroke,
        mapper.width(stroke_width * abs(item.transform.a)),
        cap=item.styles.get("stroke-linecap"),
        join=item.styles.get("stroke-linejoin"),
    )


def _requires_svg_fallback(item: RenderElement) -> bool:
    unsupported = {
        "filter",
        "mask",
        "clip-path",
        "marker",
        "marker-start",
        "marker-mid",
        "marker-end",
        "stroke-dasharray",
        "stroke-dashoffset",
    }
    return any("url(" in item.styles.get(key, "").lower() for key in ("fill", "stroke")) or any(
        item.styles.get(name) not in {None, "", "none"}
        or item.element.get(name) not in {None, "", "none"}
        for name in unsupported
    )


def _opacity(value: str | None) -> float:
    if value is None:
        return 1.0
    try:
        return max(0.0, min(1.0, float(value)))
    except ValueError:
        return 1.0


def _length(
    element: etree._Element,
    name: str,
    document: SvgDocument,
    default: float | None = None,
    *,
    horizontal: bool = True,
) -> float:
    base = document.canvas_width if horizontal else document.canvas_height
    value = parse_length(element.get(name), percentage_base=base)
    if value is None:
        if default is None:
            raise NativeShapeUnsupported(f"<{item_name(element)}> requires {name}")
        return default
    return value


def item_name(element: etree._Element) -> str:
    return element.tag.rsplit("}", 1)[-1]


def _transformed_box(
    matrix: Matrix,
    x: float,
    y: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    if width < 0 or height < 0:
        raise NativeShapeUnsupported("negative shape size")
    points = [
        matrix.apply(x, y),
        matrix.apply(x + width, y),
        matrix.apply(x, y + height),
        matrix.apply(x + width, y + height),
    ]
    xs, ys = zip(*points, strict=True)
    return min(xs), min(ys), max(xs), max(ys)
