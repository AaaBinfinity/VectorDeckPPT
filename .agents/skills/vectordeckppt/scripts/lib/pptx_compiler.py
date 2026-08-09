from __future__ import annotations

import json
from copy import deepcopy
from pathlib import Path

import resvg_py
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from .coordinates import (
    DEFAULT_SLIDE_HEIGHT_INCHES,
    DEFAULT_SLIDE_WIDTH_INCHES,
    CoordinateMapper,
)
from .path_safety import PathSafetyError, ensure_distinct_paths, require_suffix
from .pptx_images import add_native_image
from .pptx_shapes import NativeShapeUnsupported, add_native_shape
from .pptx_text import add_native_text
from .pptx_utils import (
    CompilationReport,
    SlideCompilationReport,
    add_svg_picture,
    natural_sort_key,
)
from .svg_models import RenderElement, SvgDocument
from .svg_parser import SVG_NAMESPACE, iter_render_elements, parse_svg
from .svg_validator import validate_svg

NATIVE_SHAPE_TAGS = {"rect", "circle", "ellipse", "line"}
FREEFORM_TAGS = {"polyline", "polygon"}
FALLBACK_TAGS = {"path"}


class PptxCompileError(RuntimeError):
    def __init__(self, message: str, report: CompilationReport | None = None):
        super().__init__(message)
        self.report = report


def discover_slides(source: str | Path) -> list[Path]:
    path = Path(source).expanduser().resolve()
    if path.is_file():
        return [path] if path.suffix.lower() == ".svg" else []
    if path.is_dir():
        return sorted(path.glob("*.svg"), key=natural_sort_key)
    return []


def compile_pptx(
    source: str | Path,
    output: str | Path,
) -> CompilationReport:
    source_path = Path(source).expanduser().resolve()
    output_candidate = Path(output).expanduser().resolve()
    report = CompilationReport(source=str(source_path), output=str(output_candidate))
    try:
        output_path = require_suffix(output, ".pptx", label="PowerPoint output")
    except PathSafetyError as exc:
        raise PptxCompileError(str(exc), report) from exc

    slide_paths = discover_slides(source_path)
    if not slide_paths:
        raise PptxCompileError(f"No slide SVG files found: {source_path}", report)
    try:
        for slide_path in slide_paths:
            ensure_distinct_paths(
                {"PowerPoint output": output_path, "source slide": slide_path}
            )
    except PathSafetyError as exc:
        raise PptxCompileError(str(exc), report) from exc

    presentation = Presentation()
    presentation.slide_width = Inches(DEFAULT_SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(DEFAULT_SLIDE_HEIGHT_INCHES)
    blank_layout = presentation.slide_layouts[6]

    for slide_path in slide_paths:
        slide_report = SlideCompilationReport(source=str(slide_path))
        report.slides.append(slide_report)
        validation = validate_svg(slide_path)
        if not validation.valid:
            slide_report.failed += len(validation.errors)
            slide_report.warnings.extend(
                f"{item.code}: {item.message}" for item in validation.errors
            )
            continue
        document = parse_svg(slide_path)
        view_x, view_y, view_width, view_height = document.view_box or (
            0.0,
            0.0,
            document.canvas_width,
            document.canvas_height,
        )
        mapper = CoordinateMapper(
            canvas_width=view_width,
            canvas_height=view_height,
            view_x=view_x,
            view_y=view_y,
        )
        slide = presentation.slides.add_slide(blank_layout)
        _compile_slide(slide, document, mapper, slide_report)

    if report.failed:
        raise PptxCompileError(
            f"Compilation failed with {report.failed} element or validation error(s)",
            report,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return report


def _compile_slide(
    slide: object,
    document: SvgDocument,
    mapper: CoordinateMapper,
    report: SlideCompilationReport,
) -> None:
    for item in iter_render_elements(document):
        try:
            if item.tag in NATIVE_SHAPE_TAGS:
                add_native_shape(slide, item, document, mapper)
                report.native += 1
            elif item.tag in FREEFORM_TAGS:
                add_native_shape(slide, item, document, mapper)
                report.freeform += 1
            elif item.tag == "text":
                add_native_text(slide, item, mapper)
                report.native += 1
            elif item.tag == "image":
                add_native_image(slide, item, document, mapper)
                report.native += 1
            elif item.tag in FALLBACK_TAGS:
                _add_fallback(slide, document, item, mapper)
                report.embedded_svg += 1
                report.warnings.append(
                    f"{item.xpath}: <{item.tag}> preserved as an embedded SVG fallback"
                )
            else:
                raise NativeShapeUnsupported(f"unsupported visible element <{item.tag}>")
        except NativeShapeUnsupported as exc:
            try:
                _add_fallback(slide, document, item, mapper)
                report.embedded_svg += 1
                report.warnings.append(f"{item.xpath}: {exc}; embedded SVG fallback used")
            except Exception as fallback_exc:
                report.failed += 1
                report.warnings.append(
                    f"{item.xpath}: native conversion failed ({exc}); "
                    f"SVG fallback failed ({fallback_exc})"
                )
        except Exception as exc:
            report.failed += 1
            report.warnings.append(f"{item.xpath}: unexpected compilation failure: {exc}")


def _add_fallback(
    slide: object,
    document: SvgDocument,
    item: RenderElement,
    mapper: CoordinateMapper,
) -> None:
    svg_data = _isolated_svg(document, item)
    preview = resvg_py.svg_to_bytes(
        svg_string=svg_data.decode("utf-8"),
        width=round(mapper.canvas_width),
        height=round(mapper.canvas_height),
        resources_dir=str(document.source.parent),
        shape_rendering="geometric_precision",
        text_rendering="optimize_legibility",
        image_rendering="optimize_quality",
    )
    add_svg_picture(
        slide,
        svg_data,
        preview,
        left=mapper.x(mapper.view_x),
        top=mapper.y(mapper.view_y),
        width=mapper.width(mapper.canvas_width),
        height=mapper.height(mapper.canvas_height),
        name=f"SVG fallback {item.tag}",
    )


def _isolated_svg(document: SvgDocument, item: RenderElement) -> bytes:
    nsmap = dict(document.root.nsmap)
    nsmap.setdefault(None, SVG_NAMESPACE)
    root = etree.Element(f"{{{SVG_NAMESPACE}}}svg", nsmap=nsmap)
    width = document.width or document.canvas_width
    height = document.height or document.canvas_height
    view_box = document.view_box or (0.0, 0.0, document.canvas_width, document.canvas_height)
    root.set("width", f"{width:g}")
    root.set("height", f"{height:g}")
    root.set("viewBox", " ".join(f"{value:g}" for value in view_box))

    for child in document.root:
        if isinstance(child.tag, str) and child.tag.rsplit("}", 1)[-1] == "defs":
            root.append(deepcopy(child))

    wrapper = etree.SubElement(root, f"{{{SVG_NAMESPACE}}}g")
    wrapper.set("transform", item.transform.to_svg())
    wrapper.set("opacity", f"{item.opacity:.8g}")
    clone = deepcopy(item.element)
    clone.attrib.pop("transform", None)
    clone.attrib.pop("style", None)
    clone.attrib.pop("opacity", None)
    for key, value in item.styles.items():
        if key not in {"display", "visibility", "opacity"}:
            clone.set(key, value)
    wrapper.append(clone)
    return etree.tostring(root, encoding="utf-8", xml_declaration=True)


def report_json(report: CompilationReport) -> str:
    return json.dumps(report.to_dict(), ensure_ascii=False, indent=2)
