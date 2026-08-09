from pathlib import Path

from lib.pptx_compiler import compile_pptx
from lib.pptx_validator import validate_pptx
from lib.svg_renderer import render_directory
from lib.svg_validator import validate_svg
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
EXAMPLE = ROOT / "examples" / "basic-deck"


def test_basic_deck_end_to_end(tmp_path: Path) -> None:
    slides = sorted(EXAMPLE.glob("slide_*.svg"))
    assert len(slides) == 5
    assert all(validate_svg(slide).valid for slide in slides)

    previews = render_directory(EXAMPLE, tmp_path / "preview")
    assert len(previews) == 5

    output = tmp_path / "example.pptx"
    report = compile_pptx(EXAMPLE, output)
    validation = validate_pptx(output)
    deck = Presentation(output)

    assert report.valid
    assert report.failed == 0
    assert report.embedded_svg == 0
    assert validation.valid, validation.to_dict()
    assert len(deck.slides) == 5
    assert all(any(shape.has_text_frame for shape in slide.shapes) for slide in deck.slides)
