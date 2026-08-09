from pathlib import Path

import pytest
from lib.pptx_compiler import compile_pptx
from lib.svg_validator import validate_svg
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FIXTURES = Path(__file__).parent / "fixtures"


def test_polygon_and_polyline_compile_as_editable_freeforms(tmp_path: Path) -> None:
    source = FIXTURES / "simple_freeforms.svg"
    output = tmp_path / "freeforms.pptx"

    validation = validate_svg(source)
    report = compile_pptx(source, output)
    deck = Presentation(output)
    shapes = list(deck.slides[0].shapes)

    assert validation.valid
    assert "embedded_svg_fallback" not in {item.code for item in validation.warnings}
    assert report.valid
    assert report.native == 0
    assert report.freeform == 2
    assert report.embedded_svg == 0
    assert [shape.name for shape in shapes] == ["editable-polygon", "editable-polyline"]
    assert all(shape.shape_type == MSO_SHAPE_TYPE.FREEFORM for shape in shapes)
    assert "<a:round/>" in shapes[0]._element.xml
    assert shapes[0].left == pytest.approx(1_219_200, abs=2)
    assert shapes[0].top == pytest.approx(1_371_600, abs=2)


@pytest.mark.parametrize(
    ("tag", "points", "minimum"),
    [("polyline", "10,10", 2), ("polygon", "10,10 20,20", 3)],
)
def test_freeform_validation_requires_enough_points(
    tmp_path: Path,
    tag: str,
    points: str,
    minimum: int,
) -> None:
    source = tmp_path / f"invalid-{tag}.svg"
    source.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
        f'viewBox="0 0 1600 900"><{tag} points="{points}"/></svg>',
        encoding="utf-8",
    )

    result = validate_svg(source)

    assert not result.valid
    assert any(
        item.code == "invalid_geometry" and f"at least {minimum}" in item.message
        for item in result.errors
    )
