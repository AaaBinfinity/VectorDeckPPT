from pathlib import Path

from lib.pptx_compiler import compile_pptx
from pptx import Presentation

FIXTURES = Path(__file__).parent / "fixtures"


def test_round_svg_linecap_is_preserved_in_powerpoint(tmp_path: Path) -> None:
    output = tmp_path / "round-line.pptx"

    report = compile_pptx(FIXTURES / "simple_line.svg", output)
    deck = Presentation(output)
    shape = deck.slides[0].shapes[0]

    assert report.valid
    assert 'cap="rnd"' in shape._element.xml


def test_dashed_line_uses_explicit_svg_fallback(tmp_path: Path) -> None:
    source = tmp_path / "dashed-line.svg"
    source.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
        'viewBox="0 0 1600 900"><line x1="100" y1="100" x2="500" y2="100" '
        'stroke="#2563EB" stroke-width="8" stroke-dasharray="16 8"/></svg>',
        encoding="utf-8",
    )
    output = tmp_path / "dashed-line.pptx"

    report = compile_pptx(source, output)

    assert report.valid
    assert report.native == 0
    assert report.embedded_svg == 1
    assert any("dashed stroke" in warning for warning in report.slides[0].warnings)


def test_group_inherited_dash_uses_explicit_freeform_fallback(tmp_path: Path) -> None:
    source = tmp_path / "group-dashed-polyline.svg"
    source.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
        'viewBox="0 0 1600 900"><g fill="none" stroke="#2563EB" '
        'stroke-width="8" stroke-dasharray="16 8"><polyline '
        'points="100,300 500,100 900,300"/></g></svg>',
        encoding="utf-8",
    )
    output = tmp_path / "group-dashed-polyline.pptx"

    report = compile_pptx(source, output)

    assert report.valid
    assert report.freeform == 0
    assert report.embedded_svg == 1
    assert any("dashed stroke" in warning for warning in report.slides[0].warnings)


def test_marker_definition_is_preserved_by_svg_fallback(tmp_path: Path) -> None:
    source = tmp_path / "marker-polyline.svg"
    source.write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" viewBox="0 0 1600 900">
  <defs>
    <marker id="arrow" markerWidth="10" markerHeight="10" refX="8" refY="3" orient="auto">
      <path d="M0,0 L0,6 L9,3 z" fill="#2563EB"/>
    </marker>
  </defs>
  <polyline points="100,300 500,100 900,300" fill="none" stroke="#2563EB"
            stroke-width="8" marker-end="url(#arrow)"/>
</svg>
""",
        encoding="utf-8",
    )
    output = tmp_path / "marker-polyline.pptx"

    report = compile_pptx(source, output)

    assert report.valid
    assert report.freeform == 0
    assert report.embedded_svg == 1
    assert any("marker" in warning for warning in report.slides[0].warnings)
