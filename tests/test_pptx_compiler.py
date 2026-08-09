import json
import subprocess
import sys
from pathlib import Path
from zipfile import ZipFile

import pytest
from lib.coordinates import DEFAULT_SLIDE_HEIGHT_INCHES, DEFAULT_SLIDE_WIDTH_INCHES, EMU_PER_INCH
from lib.pptx_compiler import PptxCompileError, compile_pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
FIXTURES = Path(__file__).parent / "fixtures"
COMPILER = ROOT / ".agents" / "skills" / "vectordeckppt" / "scripts" / "compile_pptx.py"


def test_compile_mixed_slide_to_editable_objects(tmp_path: Path) -> None:
    output = tmp_path / "mixed.pptx"
    report = compile_pptx(FIXTURES / "mixed_slide.svg", output)
    deck = Presentation(output)

    assert report.valid
    assert report.native >= 8
    assert report.embedded_svg == 0
    assert len(deck.slides) == 1
    assert deck.slide_width / EMU_PER_INCH == pytest.approx(DEFAULT_SLIDE_WIDTH_INCHES)
    assert deck.slide_height / EMU_PER_INCH == pytest.approx(DEFAULT_SLIDE_HEIGHT_INCHES)

    shapes = list(deck.slides[0].shapes)
    text = "\n".join(shape.text for shape in shapes if shape.has_text_frame)
    assert "VectorDeckPPT 集成测试" in text
    assert "基础图形保持可编辑" in text
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes)
    assert sum(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in shapes) >= 4
    assert all(
        "<p:style>" not in shape._element.xml
        for shape in shapes
        if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE}
    )


def test_compile_multiple_slides_uses_natural_filename_order(tmp_path: Path) -> None:
    slides = tmp_path / "slides"
    slides.mkdir()
    templates = {
        "slide_10.svg": "第十页",
        "slide_02.svg": "第二页",
        "slide_01.svg": "第一页",
    }
    for name, title in templates.items():
        (slides / name).write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
            f'viewBox="0 0 1600 900"><text x="100" y="100" font-size="40">{title}</text></svg>',
            encoding="utf-8",
        )

    output = tmp_path / "deck.pptx"
    report = compile_pptx(slides, output)
    deck = Presentation(output)
    titles = [
        next(shape.text for shape in slide.shapes if shape.has_text_frame) for slide in deck.slides
    ]

    assert report.valid
    assert titles == ["第一页", "第二页", "第十页"]


def test_complex_path_uses_office_svg_fallback(tmp_path: Path) -> None:
    output = tmp_path / "fallback.pptx"
    report = compile_pptx(FIXTURES / "fallback_path.svg", output)

    assert report.valid
    assert report.native == 1
    assert report.embedded_svg == 1
    with ZipFile(output) as archive:
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        assert any(name.endswith(".svg") for name in media)
        assert any(name.endswith(".png") for name in media)
        slide_xml = archive.read("ppt/slides/slide1.xml")
        assert b"svgBlip" in slide_xml


def test_invalid_svg_prevents_output(tmp_path: Path) -> None:
    output = tmp_path / "invalid.pptx"
    with pytest.raises(PptxCompileError) as error:
        compile_pptx(FIXTURES / "invalid_script.svg", output)

    assert error.value.report is not None
    assert error.value.report.failed > 0
    assert not output.exists()


def test_compiler_cli_writes_json_report(tmp_path: Path) -> None:
    output = tmp_path / "cli.pptx"
    report_path = tmp_path / "report.json"
    completed = subprocess.run(
        [
            sys.executable,
            str(COMPILER),
            str(FIXTURES / "simple_text.svg"),
            "--output",
            str(output),
            "--report",
            str(report_path),
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )

    assert completed.returncode == 0, completed.stderr
    payload = json.loads(report_path.read_text(encoding="utf-8"))
    assert payload["valid"] is True
    assert payload["slide_count"] == 1
    assert output.is_file()
