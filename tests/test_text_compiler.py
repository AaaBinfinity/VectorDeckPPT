from pathlib import Path

from lib.pptx_compiler import compile_pptx
from pptx import Presentation


def test_pretty_print_whitespace_does_not_create_text_runs(tmp_path: Path) -> None:
    source = tmp_path / "pretty-text.svg"
    source.write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" viewBox="0 0 1600 900">
  <text x="100" y="120" font-size="40">Alpha<tspan fill="#2563EB">Beta</tspan>
    <tspan x="100" dy="60">Gamma</tspan>
  </text>
</svg>
""",
        encoding="utf-8",
    )
    output = tmp_path / "pretty-text.pptx"

    report = compile_pptx(source, output)
    deck = Presentation(output)
    text_shapes = [shape for shape in deck.slides[0].shapes if shape.has_text_frame]
    run_texts = [
        [run.text for run in shape.text_frame.paragraphs[0].runs]
        for shape in text_shapes
    ]

    assert report.valid
    assert run_texts == [["Alpha", "Beta"], ["Gamma"]]


def test_xml_space_preserve_keeps_intentional_whitespace(tmp_path: Path) -> None:
    source = tmp_path / "preserve-text.svg"
    source.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
        'viewBox="0 0 1600 900"><text x="100" y="120" font-size="40" '
        'xml:space="preserve">Alpha<tspan>Beta</tspan> </text></svg>',
        encoding="utf-8",
    )
    output = tmp_path / "preserve-text.pptx"

    compile_pptx(source, output)
    deck = Presentation(output)
    shape = next(shape for shape in deck.slides[0].shapes if shape.has_text_frame)

    assert [run.text for run in shape.text_frame.paragraphs[0].runs] == ["Alpha", "Beta", " "]


def test_non_alphabetic_baseline_uses_explicit_svg_fallback(tmp_path: Path) -> None:
    source = tmp_path / "middle-baseline.svg"
    source.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="900" '
        'viewBox="0 0 1600 900"><text x="100" y="120" font-size="40" '
        'dominant-baseline="middle">Centered</text></svg>',
        encoding="utf-8",
    )
    output = tmp_path / "middle-baseline.pptx"

    report = compile_pptx(source, output)

    assert report.valid
    assert report.native == 0
    assert report.embedded_svg == 1
    assert any("dominant-baseline" in warning for warning in report.slides[0].warnings)
